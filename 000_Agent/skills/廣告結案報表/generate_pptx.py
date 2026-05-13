#!/usr/bin/env python3
"""
廣告結案報表產生器（模板修改版）
Usage: python3 generate_pptx.py <data.json> <template.pptx> [output.pptx]

直接修改既有 PPTX 模板中的數字與文字，保留所有版型設計。
"""
import json, sys, re
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
         'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}


# ── XML helpers ───────────────────────────────────────────────────────────────

def shape_text(shape):
    """全部段落文字串接。"""
    if not shape.has_text_frame:
        return ""
    return " ".join("".join(r.text for r in p.runs)
                    for p in shape.text_frame.paragraphs).strip()


def set_run_text(run, text):
    run.text = text


def clear_and_set_tf(tf, lines):
    """
    lines: list of (text, size_pt, bold, color_hex_or_None)
    清空 text frame 後依序填入每行。
    """
    txBody = tf._txBody
    # 保留第一個 <a:p> 的 pPr（對齊等屬性），刪掉其他段落
    paras = txBody.findall(qn('a:p'))
    # 取第一個段落的 pPr 作為格式模板
    first_pPr = None
    if paras:
        first_pPr = deepcopy(paras[0].find(qn('a:pPr')))
    # 刪掉所有段落
    for p in paras:
        txBody.remove(p)

    for i, (text, size_pt, bold, color_hex) in enumerate(lines):
        p_el = etree.SubElement(txBody, qn('a:p'))
        if first_pPr is not None and i == 0:
            p_el.insert(0, deepcopy(first_pPr))
        r_el = etree.SubElement(p_el, qn('a:r'))
        rPr = etree.SubElement(r_el, qn('a:rPr'))
        rPr.set('lang', 'zh-TW')
        if size_pt:
            rPr.set('sz', str(int(size_pt * 100)))
        if bold:
            rPr.set('b', '1')
        if color_hex:
            sf = etree.SubElement(rPr, qn('a:solidFill'))
            srgb = etree.SubElement(sf, qn('a:srgbClr'))
            srgb.set('val', color_hex.upper())
        t_el = etree.SubElement(r_el, qn('a:t'))
        t_el.text = text


def find_shape_by_keyword(slide, *keywords):
    """找到包含任一 keyword 的 shape。"""
    for shp in slide.shapes:
        t = shape_text(shp)
        if any(kw in t for kw in keywords):
            return shp
    return None


def find_table(slide):
    for shp in slide.shapes:
        if shp.has_table:
            return shp
    return None


def find_chart(slide):
    for shp in slide.shapes:
        if shp.has_chart:
            return shp
    return None


def find_all_charts(slide):
    return [shp for shp in slide.shapes if shp.has_chart]


# ── Per-slide updaters ────────────────────────────────────────────────────────

def update_cover(slide, d):
    """Slide 1: 封面 — 更新走期 & 月份。"""
    period_shp = find_shape_by_keyword(slide, '專案走期', '走期')
    if period_shp:
        clear_and_set_tf(period_shp.text_frame,
                         [(f'專案走期：{d["period"]}', 32, True, None)])

    title_shp = find_shape_by_keyword(slide, '結案報告', '廣告專案')
    if title_shp:
        clear_and_set_tf(title_shp.text_frame,
                         [(f'{d["client"]} 廣告專案 – {d["month"]} 月結案報告', 60, True, None)])


def update_kpi(slide, d):
    """Slide 2: 成效總覽 — 替換四個 KPI 數字。"""
    kpi = d['kpi']
    period = d['period'].split('~')[0].replace('/', '/').strip()
    end = d['period'].split('~')[1].strip()
    short_period = f'({period}~{end})'

    for shp in slide.shapes:
        t = shape_text(shp)
        if '總花費' in t:
            clear_and_set_tf(shp.text_frame, [
                (f'總花費', 29, True, None),
                (short_period, 20, True, None),
                (f'NT${int(kpi["total_spend"]):,}', 36, True, None),
            ])
        elif '總曝光數' in t:
            clear_and_set_tf(shp.text_frame, [
                ('總曝光數', 29, True, None),
                (f'{int(kpi["total_impressions"]):,}', 36, True, None),
                ('次', 25, True, None),
            ])
        elif '總點擊數' in t:
            clear_and_set_tf(shp.text_frame, [
                ('總點擊數', 29, True, None),
                (f'{int(kpi["total_clicks"]):,}', 36, True, None),
                ('次', 25, True, None),
            ])
        elif '總購買次數' in t:
            clear_and_set_tf(shp.text_frame, [
                ('總購買次數', 29, True, None),
                (f'{int(kpi["total_purchases"]):,}', 36, True, None),
                ('筆', 25, True, None),
            ])


def update_impressions(slide, d):
    """Slide 3: 曝光及點擊狀況。"""
    kpi = d['kpi']
    period_str = d['period'].replace('~', '–')
    summary = (f'走期 {period_str}，累積曝光 {int(kpi["total_impressions"]):,} 次，'
               f'點擊 {int(kpi["total_clicks"]):,} 次。')
    shp = find_shape_by_keyword(slide, '累積曝光', '走期', '點擊狀況')
    if shp and '曝光及點擊狀況' not in shape_text(shp):  # 避免改到標題
        clear_and_set_tf(shp.text_frame, [(summary, 25, False, None)])

    daily = d.get('daily_data', {})
    chart_shp = find_chart(slide)
    if chart_shp and daily.get('dates'):
        cd = ChartData()
        cd.categories = daily['dates']
        cd.add_series('曝光數', daily['impressions'])
        cd.add_series('點擊數', daily['clicks'])
        chart_shp.chart.replace_data(cd)


def update_purchase(slide, d):
    """Slide 4: 購買次數狀況。"""
    kpi = d['kpi']
    summary = (f'廣告累積購買數 {int(kpi["total_purchases"]):,} 次，'
               f'購買轉換值 ${int(kpi["total_purchase_value"]):,}，'
               f'廣告投報率為 {kpi["roas"]:.2f}。')
    shp = find_shape_by_keyword(slide, '累積購買', '購買轉換')
    if shp:
        clear_and_set_tf(shp.text_frame, [(summary, 25, False, None)])

    daily = d.get('daily_data', {})
    chart_shp = find_chart(slide)
    if chart_shp and daily.get('purchases'):
        cd = ChartData()
        cd.categories = daily['dates']
        cd.add_series('購買數', daily['purchases'])
        chart_shp.chart.replace_data(cd)


def update_table(slide, d):
    """Slide 5: 素材分析表格。"""
    tbl_shp = find_table(slide)
    if tbl_shp is None:
        return

    table = tbl_shp.table
    creatives = d['creatives']
    kpi = d['kpi']

    # 取得表頭（保留原始格式）
    # 先記錄每個欄位的寬度
    header_row = table.rows[0]

    # 刪除舊資料列（保留表頭），Python-pptx 無法直接刪列，用 XML 操作
    tbl_xml = tbl_shp.table._tbl
    rows_xml = tbl_xml.findall(qn('a:tr'))

    # 刪除資料列（index 1 以後），保留表頭（index 0）
    for r in rows_xml[1:]:
        tbl_xml.remove(r)

    # 取得表頭第一列的 XML 作為格式模板
    header_tr = rows_xml[0]
    # 從第一個資料列（如果存在）取格式；否則從表頭推斷
    template_tr = rows_xml[1] if len(rows_xml) > 1 else None

    def make_cell(text, bold=False, color_hex=None, align='center'):
        """產生 <a:tc> XML element。"""
        tc = etree.Element(qn('a:tc'))
        txBody = etree.SubElement(tc, qn('a:txBody'))
        bodyPr = etree.SubElement(txBody, qn('a:bodyPr'))
        lstStyle = etree.SubElement(txBody, qn('a:lstStyle'))
        p = etree.SubElement(txBody, qn('a:p'))
        pPr = etree.SubElement(p, qn('a:pPr'))
        if align == 'center':
            pPr.set('algn', 'ctr')
        elif align == 'left':
            pPr.set('algn', 'l')
        r = etree.SubElement(p, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'))
        rPr.set('lang', 'zh-TW')
        rPr.set('sz', '800')
        if bold:
            rPr.set('b', '1')
        if color_hex:
            sf = etree.SubElement(rPr, qn('a:solidFill'))
            srgb = etree.SubElement(sf, qn('a:srgbClr'))
            srgb.set('val', color_hex)
        t = etree.SubElement(r, qn('a:t'))
        t.text = text
        # 加 tcPr（邊框）
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
        return tc

    def make_row(cells_data, is_total=False):
        tr = etree.Element(qn('a:tr'))
        tr.set('h', '400000')  # 約 0.44 inch 行高
        for text, align in cells_data:
            bold = is_total
            tc = make_cell(text, bold=bold, align=align)
            tr.append(tc)
        return tr

    # 資料列定義
    def creative_row(c):
        lcc = f'${c.get("link_click_cost", 0):.2f}' if c.get("link_click_cost") else '-'
        return make_row([
            (c.get('format', ''), 'center'),
            (c.get('name', ''), 'left'),
            (f'{int(c.get("impressions", 0)):,}', 'center'),
            (f'{int(c.get("clicks", 0)):,}', 'center'),
            (f'{c.get("ctr", 0):.2f}%', 'center'),
            (f'{int(c.get("video_views", 0)):,}', 'center'),
            (f'{int(c.get("link_clicks", 0)):,}', 'center'),
            (lcc, 'center'),
            (f'{int(c.get("purchases", 0)):,}', 'center'),
            (f'${int(c.get("purchase_value", 0)):,}', 'center'),
            (f'{c.get("roas", 0):.2f}', 'center'),
            (f'${int(c.get("spend", 0)):,}', 'center'),
        ])

    total_ctr = kpi['total_clicks'] / kpi['total_impressions'] * 100 if kpi['total_impressions'] else 0
    total_row = make_row([
        ('總計', 'center'),
        ('', 'center'),
        (f'{int(kpi["total_impressions"]):,}', 'center'),
        (f'{int(kpi["total_clicks"]):,}', 'center'),
        (f'{total_ctr:.2f}%', 'center'),
        ('', 'center'),
        ('', 'center'),
        (f'${kpi["total_spend"]/kpi["total_clicks"]:.2f}' if kpi["total_clicks"] else '-', 'center'),
        (f'{int(kpi["total_purchases"]):,}', 'center'),
        (f'${int(kpi["total_purchase_value"]):,}', 'center'),
        (f'{kpi["roas"]:.2f}', 'center'),
        (f'${int(kpi["total_spend"]):,}', 'center'),
    ], is_total=True)

    for c in creatives:
        tbl_xml.append(creative_row(c))
    tbl_xml.append(total_row)

    # 更新 summary 文字
    summary = d.get('creative_summary', '')
    if summary:
        shp = find_shape_by_keyword(slide, '整體曝光', '表現', '廣告投遞')
        if shp and not shp.has_table:
            clear_and_set_tf(shp.text_frame, [(summary, 23, False, None)])


def update_screenshots(slides, d):
    """Slides 6-8: 廣告上刊截圖頁 — 更新素材名稱與指標文字。"""
    creatives = d['creatives']
    # 每頁最多 3 則
    batches = [creatives[i:i+3] for i in range(0, len(creatives), 3)]

    for slide_idx, (slide, batch) in enumerate(zip(slides, batches)):
        # 找出所有「黑色標籤框」（fill=000000，含素材名稱）
        # 和指標文字框（含「曝光數」「點擊數」）
        name_shapes = []
        metric_shapes = []
        for shp in slide.shapes:
            t = shape_text(shp)
            if shp.has_text_frame:
                # 嘗試讀 fill color
                try:
                    fc = shp.fill.fore_color.rgb
                    if str(fc) == '000000':
                        name_shapes.append(shp)
                        continue
                except Exception:
                    pass
                if '曝光數' in t and '點擊數' in t:
                    metric_shapes.append(shp)

        # 依左到右排序
        name_shapes.sort(key=lambda s: s.left)
        metric_shapes.sort(key=lambda s: s.left)

        for i, c in enumerate(batch):
            if i < len(name_shapes):
                clear_and_set_tf(name_shapes[i].text_frame, [(c['name'], 20, True, '000000')])
            if i < len(metric_shapes):
                vv = c.get('video_views', 0)
                vv_str = f'約 {vv/10000:.1f} 萬次' if vv >= 10000 else f'{vv:,} 次'
                lines = [
                    (f'曝光數：{int(c.get("impressions", 0)):,}', 16, False, None),
                    (f'點擊數：{int(c.get("clicks", 0)):,}', 16, False, None),
                    (f'影片觀看：{vv_str}', 16, False, None),
                    (f'連結點擊數：{int(c.get("link_clicks", 0)):,}', 16, False, None),
                ]
                clear_and_set_tf(metric_shapes[i].text_frame, lines)


def update_audience(slide, d):
    """Slide 9: 受眾分析 — 更新分析文字與圖表。"""
    audience = d.get('audience', {})
    analysis = audience.get('analysis', '')

    shp = find_shape_by_keyword(slide, '購買年齡', '購買者', '性別比例', '女性族群')
    if shp:
        clear_and_set_tf(shp.text_frame, [(analysis, 24, False, None)])

    # 更新年齡圖表
    age_groups = audience.get('age_groups', [])
    gender = audience.get('gender', {})
    charts = find_all_charts(slide)

    if charts and age_groups:
        cd = ChartData()
        cd.categories = [g['range'] for g in age_groups]
        cd.add_series('購買占比(%)', [g['percentage'] for g in age_groups])
        try:
            charts[0].chart.replace_data(cd)
        except Exception:
            pass

    if len(charts) > 1 and gender:
        cd2 = ChartData()
        cd2.categories = ['女性', '男性']
        cd2.add_series('占比', [gender.get('female', 0), gender.get('male', 0)])
        try:
            charts[1].chart.replace_data(cd2)
        except Exception:
            pass


def update_conclusion(slide, d):
    """Slide 10: 整體廣告概述 + 分類建議。"""
    overview = d.get('conclusion', '')

    # 左欄：整體廣告概述
    for shp in slide.shapes:
        t = shape_text(shp)
        if ('本次廣告' in t or '整體曝光' in t or '購買訂單' in t) and '廣告結論' not in t:
            clear_and_set_tf(shp.text_frame, [(overview, 21, False, None)])
            break

    # 右欄：分類建議（素材建議 / 文案建議 / 活動建議）
    creative_recs  = d.get('creative_recommendations', [])
    copy_recs      = d.get('copy_recommendations', [])
    campaign_recs  = d.get('campaign_recommendations', [])
    # 相容舊版 recommendations（無分類）
    fallback_recs  = d.get('recommendations', [])

    lines = []
    if creative_recs or copy_recs or campaign_recs:
        if creative_recs:
            lines.append(('【素材建議】', 21, True, None))
            for r in creative_recs:
                lines.append((f'• {r}', 21, False, None))
        if copy_recs:
            if lines: lines.append(('', 10, False, None))
            lines.append(('【文案建議】', 21, True, None))
            for r in copy_recs:
                lines.append((f'• {r}', 21, False, None))
        if campaign_recs:
            if lines: lines.append(('', 10, False, None))
            lines.append(('【活動建議】', 21, True, None))
            for r in campaign_recs:
                lines.append((f'• {r}', 21, False, None))
    else:
        lines = [(f'• {r}', 21, False, None) for r in fallback_recs]

    for shp in slide.shapes:
        t = shape_text(shp)
        if '建議方向' in t and len(t) > 10:
            clear_and_set_tf(shp.text_frame, [('建議方向', 21, True, None)] + lines)
            break


# ── Main ──────────────────────────────────────────────────────────────────────

def generate(data_path, template_path, output_path=None):
    with open(data_path, encoding='utf-8') as f:
        d = json.load(f)

    prs = Presentation(template_path)
    slides = list(prs.slides)
    n = len(slides)

    if n >= 1:  update_cover(slides[0], d)
    if n >= 2:  update_kpi(slides[1], d)
    if n >= 3:  update_impressions(slides[2], d)
    if n >= 4:  update_purchase(slides[3], d)
    if n >= 5:  update_table(slides[4], d)
    if n >= 8:  update_screenshots(slides[5:8], d)
    if n >= 9:  update_audience(slides[8], d)
    if n >= 10: update_conclusion(slides[9], d)
    # Slide 11 (Thanks) — 不動

    if not output_path:
        client = d.get('client', 'client')
        month = d.get('month', '')
        output_path = f"{client}_{month}月結案報表.pptx"

    prs.save(output_path)
    print(f"✅ 報表已產出：{output_path}")
    return output_path


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python3 generate_pptx.py <data.json> <template.pptx> [output.pptx]")
        sys.exit(1)
    generate(sys.argv[1], sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else None)
